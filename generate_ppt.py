import json
import pptx
import random

class generate_ppt:

    def __init__(self, topic: str, category: str) -> None:
        self.topic = topic
        self.category = category.replace(' ','_')
        with open('json_format.json','r') as f:
            self.json_data = json.load(f)


    def select_random_template(self):
        self.random_category = random.choice(list(self.json_data[self.category].keys()))
        print(self.random_category)
        path = getattr(self, f'{self.category}_{self.random_category}')
        self.presentation = pptx.Presentation(f'templates/{self.category}/{self.random_category}.pptx')
        return path()
        
    def get_data(self):
        import google.generativeai as genai
        from dotenv import load_dotenv
        import os
        import json

        load_dotenv()

        genai.configure(api_key=os.getenv('GEMINI_API_KEY'))
        model = genai.GenerativeModel('gemini-pro')


        prompt = f'Fill the below json with a detailed overview of the topic: "{self.topic}"\n'

        
        json_format = json.dumps(self.json_data[self.category][self.random_category])

        prompt = prompt + json_format

        while True:
            try: 
                response = model.generate_content(prompt)

                json_response: json = json.loads(response.text)

                with open('slides.json','w') as f:
                    json.dump(json_response, f, indent=4)

                return json_response
            except: 
                print('Error Trying again...')
            
        


    def list_text_boxes(self, presentation, slide_num):
        slide = presentation.slides[slide_num-1]
        text_boxes = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                text_boxes.append(shape.text)
        return text_boxes

    def get_text_box_id(self, slide_no):
        for idx, text in enumerate(self.list_text_boxes(self.presentation, slide_no), 1):
            print(f"slide no. {slide_no}: Text Box {idx}: {text}")


    def update_text_of_textbox(self, slide, text_box_id, new_text):
        slide = self.presentation.slides[(slide - 1)]
        count = 0
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text:
                count += 1
                if count == text_box_id:
                    text_frame = shape.text_frame
                    first_paragraph = text_frame.paragraphs[0]
                    first_run = first_paragraph.runs[0] if first_paragraph.runs else first_paragraph.add_run()
                    first_run.text = new_text
                    return
        raise ValueError(f"Slide {slide} or Text Box ID {text_box_id} not found")

    def stud_project_modern(self):
        data = self.get_data()
        self.update_text_of_textbox(1,2,data['topic']) # slide 1
        
        self.update_text_of_textbox(2,2,data['introduction']) # slide 2
        
        self.update_text_of_textbox(4,2,data['background']) # slide 4
        
        self.update_text_of_textbox(6,2,data['objectives']) # slide 6

        self.update_text_of_textbox(8,1,data['methodology'][0]) # slide 8: meth 1
        self.update_text_of_textbox(8,2,data['methodology'][1]) # slide 8: meth 2
        self.update_text_of_textbox(8,3,data['methodology'][2]) # slide 8: meth 3

        self.update_text_of_textbox(9,2,data['conclusion']) # slide 9

        file_name = f"{self.topic.replace(' ','_')}.pptx"

        self.presentation.save(file_name)

        return f'PPT saved: {file_name}'
        
    def stud_project_minimalist(self):
        data = self.get_data()

        self.update_text_of_textbox(1, 3,data['topic']) # slide 1
        
        self.update_text_of_textbox(3,2,data['introduction']) # slide 3
        
        self.update_text_of_textbox(4,3,data['objectives']) # slide 4
        
        self.update_text_of_textbox(5,1,data['scope']) # slide 5

        self.update_text_of_textbox(6,3,data['background']) # slide 6

        self.update_text_of_textbox(7,2,data['methodology']) # slide 7

        self.update_text_of_textbox(8,2,data['benefit']) # slide 8

        file_name = f"{self.topic.replace(' ','_')}.pptx"

        self.presentation.save(file_name)

        return f'PPT saved: {file_name}'


if __name__ == '__main__':
    a = generate_ppt('hotal management system', 'stud_project')
    a.select_random_template()
    # a.get_text_box_id(1)
    # a.get_text_box_id(3)
    # a.get_text_box_id(4)
    # a.get_text_box_id(5)
    # a.get_text_box_id(6)
    # a.get_text_box_id(7)
    # a.get_text_box_id(8)
    # a.get_data()
    # s = a.minimalist_stud_project()
    # print(s)
